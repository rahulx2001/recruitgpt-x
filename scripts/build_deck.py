#!/usr/bin/env python3
"""Build hackathon approach deck (PPTX) for RecruitGPT X."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "RecruitGPT_X_Approach.pptx"

NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
TEAL = RGBColor(0x02, 0xC3, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x36, 0x45, 0x4F)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)


def _bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _title_box(slide, text: str, *, dark: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.8), Inches(0.9))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark else NAVY


def _bullets(slide, items: list[str], y: float = 1.5) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(y), Inches(8.5), Inches(3.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(17)
        p.font.color.rgb = CHARCOAL
        p.space_after = Pt(10)


def _accent_bar(slide) -> None:
    shape = slide.shapes.add_shape(
        1, Inches(0.55), Inches(0.35), Inches(0.08), Inches(0.75)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()


def _stat_cards(slide, stats: list[tuple[str, str]]) -> None:
    x_positions = [0.7, 3.55, 6.4]
    for (value, label), x in zip(stats, x_positions):
        card = slide.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(2.5), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = ICE
        card.line.color.rgb = NAVY
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.text = value
        p1.font.size = Pt(30)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1 — Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(1.2))
    t.text_frame.text = "RecruitGPT X"
    t.text_frame.paragraphs[0].font.size = Pt(44)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    sub = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.2), Inches(1.0))
    sub.text_frame.text = "Ranking candidates like a great recruiter — not a keyword filter"
    sub.text_frame.paragraphs[0].font.size = Pt(20)
    sub.text_frame.paragraphs[0].font.color.rgb = ICE
    tag = s.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.0), Inches(0.5))
    tag.text_frame.text = "India Runs Data & AI Challenge · Redrob Senior AI Engineer"
    tag.text_frame.paragraphs[0].font.size = Pt(14)
    tag.text_frame.paragraphs[0].font.color.rgb = TEAL

    # 2 — Problem
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _title_box(s, "The Problem")
    _bullets(
        s,
        [
            "Recruiters review hundreds of profiles but keyword filters miss real fit.",
            "Honeypot profiles stuff AI skills (HTML, Tailwind) into unrelated titles.",
            "Behavioral signals — response rate, GitHub, saves — are ignored by ATS.",
            "We need ranking that understands the role, not just the vocabulary.",
        ],
    )

    # 3 — Role understanding
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _title_box(s, "What the Role Actually Needs")
    _bullets(
        s,
        [
            "Parsed job_description.docx into structured requirements (jd_config.py).",
            "Core IR/ML stack: embeddings, retrieval, ranking, vector DBs, PyTorch.",
            "5–9 years experience, India location, production deployment signals.",
            "Redrob behavioral modifier: recruiter response, GitHub, saves, completeness.",
            "Explicit honeypot title list + noise-skill penalty for keyword stuffers.",
        ],
    )

    # 4 — Architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _title_box(s, "Hybrid Offline Ranker")
    _bullets(
        s,
        [
            "Two-layer system: offline rank.py for 100K submission + live 7-agent app.",
            "Weighted components: title 32%, skills 28%, production 18%, experience 12%, location 10%.",
            "Multiplicative modifiers: consulting-only penalty, behavioral boost, honeypot trap.",
            "Skill trust = endorsements + duration; core vs secondary vs noise buckets.",
            "CPU-only, no network — ranks 100K candidates in ~7 seconds.",
        ],
        y=1.35,
    )

    # 5 — Scoring diagram (stats)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _title_box(s, "Scoring at a Glance")
    _stat_cards(
        s,
        [
            ("32%", "Title + career alignment"),
            ("28%", "Core AI/IR skills"),
            ("18%", "Production signals"),
        ],
    )
    note = s.shapes.add_textbox(Inches(0.75), Inches(4.0), Inches(8.5), Inches(1.0))
    note.text_frame.text = (
        "Tie-breakers: raw score → core skill count → production depth → behavioral signals → ID"
    )
    note.text_frame.paragraphs[0].font.size = Pt(14)
    note.text_frame.paragraphs[0].font.color.rgb = MUTED

    # 6 — Honeypots
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _title_box(s, "Beating Keyword Traps")
    _bullets(
        s,
        [
            "Weak titles (HR Manager, Accountant, etc.) capped at 0.12 title score.",
            "Profiles with noise skills (HTML, Photoshop) get multiplicative penalty.",
            "Honeypot modifier: unrelated title + inflated AI skills → 0.25× demotion.",
            "Result: zero honeypot titles in our top-100 shortlist.",
            "Top rank: Senior AI Engineer with 7 core IR skills and 0.80 response rate.",
        ],
    )

    # 7 — Live platform
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _title_box(s, "Beyond the CSV: RecruitGPT X Platform")
    _bullets(
        s,
        [
            "Next.js dashboard + FastAPI + LangGraph 7-agent pipeline for recruiters.",
            "Semantic matching (BGE embeddings), explainability, bias audit, what-if analysis.",
            "Candidate Radar visualization, skill evolution timeline, AI recruiter chat.",
            "Same philosophy online: understand fit, explain every rank, surface hidden talent.",
        ],
    )

    # 8 — Results
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _title_box(s, "Submission Results")
    _stat_cards(
        s,
        [
            ("100K", "Candidates ranked"),
            ("7s", "CPU runtime"),
            ("0", "Honeypots in top 100"),
        ],
    )
    res = s.shapes.add_textbox(Inches(0.75), Inches(3.9), Inches(8.5), Inches(1.2))
    res.text_frame.text = (
        "#1 CAND_0002025 Senior AI Engineer · #2 CAND_0071974 Senior AI Engineer (8 core skills)\n"
        "Reproduce: python rank.py --candidates ./data/candidates.jsonl --out ./submission.csv"
    )
    for p in res.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.color.rgb = CHARCOAL

    # 9 — Why this approach
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)
    _title_box(s, "Why We Built It This Way")
    _bullets(
        s,
        [
            "Interpretable weights — judges and recruiters can audit every decision.",
            "Domain-informed features beat generic embedding-only search on honeypots.",
            "Behavioral signals used as modifier, not override — skills still matter.",
            "Reproducible, deterministic, hackathon-compliant (no GPU, no network).",
            "Live app proves the same ideas scale to interactive recruiter workflows.",
        ],
    )

    # 10 — Close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.0))
    t.text_frame.text = "Make hiring smarter."
    t.text_frame.paragraphs[0].font.size = Pt(40)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    c = s.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.2), Inches(1.2))
    c.text_frame.text = "GitHub: recruitgpt-x  ·  submission.csv validated  ·  Questions welcome"
    c.text_frame.paragraphs[0].font.size = Pt(18)
    c.text_frame.paragraphs[0].font.color.rgb = ICE

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")