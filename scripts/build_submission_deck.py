#!/usr/bin/env python3
"""Build Redrob Idea Submission deck — consultant-grade light green theme."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_PPTX = ROOT / "docs" / "Schadn_Redrob_Submission.pptx"

# Light green consultant palette
FOREST = RGBColor(0x1B, 0x43, 0x32)
SAGE = RGBColor(0x2D, 0x6A, 0x4F)
LEAF = RGBColor(0x40, 0x91, 0x6C)
MINT = RGBColor(0x74, 0xC6, 0x9D)
MINT_PALE = RGBColor(0xB7, 0xE4, 0xC7)
CREAM = RGBColor(0xF7, 0xFB, 0xF8)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
CARD_TINT = RGBColor(0xED, 0xF7, 0xF0)
INK = RGBColor(0x1A, 0x2E, 0x22)
MUTED = RGBColor(0x5C, 0x6F, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD_LINE = RGBColor(0x95, 0xD5, 0xB2)

FONT = "Calibri"
SLIDE_H = Inches(5.625)
SLIDE_W = Inches(10)
_slide_no = 0


def _blank(prs: Presentation):
    global _slide_no
    _slide_no += 1
    return prs.slides.add_slide(prs.slide_layouts[6])


def _shape(slide, kind, left, top, w, h, fill: RGBColor, *, line=None, line_w=0.75):
    s = slide.shapes.add_shape(kind, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s


def _rect(slide, left, top, w, h, fill, **kw):
    return _shape(slide, MSO_SHAPE.RECTANGLE, left, top, w, h, fill, **kw)


def _round_rect(slide, left, top, w, h, fill, **kw):
    return _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, fill, **kw)


def _text(slide, left, top, w, h, text, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color
    p.alignment = align
    return box


def _rich_bullets(slide, left, top, w, h, title, items, *, title_size=17, body_size=11.5, accent=FOREST):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = accent
    p.space_after = Pt(8)
    for item in items:
        bp = tf.add_paragraph()
        bp.text = item
        bp.level = 0
        bp.font.size = Pt(body_size)
        bp.font.name = FONT
        bp.font.color.rgb = INK
        bp.space_after = Pt(5)
        bp.line_spacing = 1.15
    return box


def _slide_canvas(slide, section: str, n: int):
    """Consultant base: cream bg, left accent rail, footer."""
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CREAM)
    _rect(slide, Inches(0), Inches(0), Inches(0.14), SLIDE_H, SAGE)
    _rect(slide, Inches(0.14), Inches(0), Inches(0.06), SLIDE_H, MINT_PALE)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.425), FOREST)
    _text(slide, Inches(0.55), Inches(5.28), Inches(5), Inches(0.28),
          "Schadn  ·  Rahul Kumar Singh  ·  India Runs Data & AI Challenge", size=8.5, color=MINT_PALE)
    _text(slide, Inches(8.6), Inches(5.28), Inches(1.2), Inches(0.28), f"{n:02d}", size=9, color=MINT, align=PP_ALIGN.RIGHT)
    if section:
        _text(slide, Inches(0.55), Inches(0.32), Inches(3), Inches(0.22), section.upper(), size=8, bold=True, color=LEAF)


def _slide_title(slide, title: str, subtitle: str = ""):
    _text(slide, Inches(0.55), Inches(0.52), Inches(8.8), Inches(0.55), title, size=28, bold=True, color=FOREST)
    if subtitle:
        _text(slide, Inches(0.55), Inches(1.02), Inches(8.8), Inches(0.32), subtitle, size=12, color=MUTED)


def _card(slide, left, top, w, h, *, tint=False):
    return _round_rect(slide, left, top, w, h, CARD_TINT if tint else CARD, line=MINT_PALE, line_w=1)


def _stat_tile(slide, left, top, value, label, sub=""):
    _round_rect(slide, left, top, Inches(2.55), Inches(1.05), CARD, line=MINT, line_w=1.2)
    _rect(slide, left, top, Inches(2.55), Inches(0.06), LEAF)
    _text(slide, left + Inches(0.18), top + Inches(0.14), Inches(2.2), Inches(0.45), value, size=26, bold=True, color=FOREST)
    _text(slide, left + Inches(0.18), top + Inches(0.58), Inches(2.2), Inches(0.22), label, size=10, color=MUTED)
    if sub:
        _text(slide, left + Inches(0.18), top + Inches(0.78), Inches(2.2), Inches(0.18), sub, size=8.5, color=LEAF)


def slide_cover(prs: Presentation):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), Inches(3.65), SLIDE_H, FOREST)
    _rect(slide, Inches(3.65), Inches(0), Inches(6.35), SLIDE_H, CREAM)
    _rect(slide, Inches(3.65), Inches(0), Inches(0.08), SLIDE_H, LEAF)
    # decorative circles
    _shape(slide, MSO_SHAPE.OVAL, Inches(0.35), Inches(4.1), Inches(2.8), Inches(2.8), RGBColor(0x23, 0x5C, 0x45))
    _shape(slide, MSO_SHAPE.OVAL, Inches(1.1), Inches(0.2), Inches(1.4), Inches(1.4), RGBColor(0x2D, 0x6A, 0x4F))

    _text(slide, Inches(0.45), Inches(0.55), Inches(3.0), Inches(0.35), "SCHADN", size=13, bold=True, color=MINT)
    _text(slide, Inches(0.45), Inches(1.0), Inches(3.0), Inches(1.4),
          "Intelligent\nCandidate\nDiscovery", size=30, bold=True, color=WHITE)
    _text(slide, Inches(0.45), Inches(2.55), Inches(3.0), Inches(0.5), "& Ranking System", size=16, color=MINT_PALE)
    _text(slide, Inches(0.45), Inches(3.35), Inches(3.0), Inches(0.7),
          "India Runs\nData & AI Challenge", size=11, color=RGBColor(0xD8, 0xF3, 0xDC))

    _text(slide, Inches(4.0), Inches(0.75), Inches(5.7), Inches(0.3), "EXECUTIVE SUBMISSION", size=9, bold=True, color=LEAF)
    _text(slide, Inches(4.0), Inches(1.05), Inches(5.7), Inches(0.9),
          "Hybrid Offline Ranker\nfor Redrob AI", size=34, bold=True, color=FOREST)
    _text(slide, Inches(4.0), Inches(2.05), Inches(5.7), Inches(0.55),
          "Semantic fit over keywords · Trap-aware · Production-scale CPU ranking", size=13, color=MUTED)

    _round_rect(slide, Inches(4.0), Inches(2.85), Inches(5.55), Inches(1.15), CARD, line=MINT_PALE)
    _text(slide, Inches(4.25), Inches(3.0), Inches(2.4), Inches(0.85), "Team\nSchadn", size=14, bold=True, color=FOREST)
    _text(slide, Inches(6.55), Inches(3.0), Inches(2.8), Inches(0.85), "Team Leader\nRahul Kumar Singh", size=14, bold=True, color=FOREST)

    _text(slide, Inches(4.0), Inches(4.2), Inches(5.55), Inches(0.75),
          "Official bundle: candidates.jsonl (100K) · job_description.docx ·\n"
          "redrob_signals_doc.docx · submission_spec v4",
          size=9.5, color=MUTED)
    _rect(slide, Inches(4.0), Inches(5.05), Inches(2.2), Inches(0.04), LEAF)


def slide_solution(prs: Presentation):
    slide = _blank(prs)
    _slide_canvas(slide, "Solution", 2)
    _slide_title(slide, "Solution Overview", "A consultant-grade ranking engine built for real recruiting constraints")

    _card(slide, Inches(0.55), Inches(1.35), Inches(4.35), Inches(3.65))
    _rect(slide, Inches(0.55), Inches(1.35), Inches(4.35), Inches(0.42), CARD_TINT)
    _text(slide, Inches(0.75), Inches(1.42), Inches(3.9), Inches(0.3), "Proposed Solution", size=15, bold=True, color=FOREST)
    _rich_bullets(slide, Inches(0.75), Inches(1.9), Inches(3.95), Inches(2.9), "", [
        "Rank top 100 from official candidates.jsonl per submission_spec v4.",
        "Reads candidate_schema: profile, career_history, skills, redrob_signals (23 fields).",
        "Hybrid offline ranker: ~49–60s on CPU — within 5 min / 16GB / no-network limits.",
        "Grounded 2-sentence reasoning per row aligned to Stage 4 review criteria.",
    ], title_size=1, body_size=11)

    _card(slide, Inches(5.1), Inches(1.35), Inches(4.35), Inches(3.65), tint=True)
    _rect(slide, Inches(5.1), Inches(1.35), Inches(4.35), Inches(0.42), MINT_PALE)
    _text(slide, Inches(5.3), Inches(1.42), Inches(3.9), Inches(0.3), "Strategic Differentiation", size=15, bold=True, color=FOREST)
    _rich_bullets(slide, Inches(5.3), Inches(1.9), Inches(3.95), Inches(2.9), "", [
        "Defends bundle traps: keyword stuffers, behavioral twins, ~80 honeypots.",
        "Career semantics + IR skill trust + Redrob behavioral modifier.",
        "Mirrors JD disqualifiers: research-only, consulting-only, CV/speech-only.",
        "open_to_work_flag=false cannot reach top-10.",
        "Shipper-first: scales to 200K+ without per-candidate LLM cost.",
    ], title_size=1, body_size=10.5)


def slide_jd(prs: Presentation):
    slide = _blank(prs)
    _slide_canvas(slide, "JD & Signals", 3)
    _slide_title(slide, "JD Understanding & Candidate Signals",
                 "Grounded in job_description.docx · candidate_schema.json · redrob_signals_doc.docx")

    _card(slide, Inches(0.55), Inches(1.28), Inches(4.4), Inches(1.95))
    _rich_bullets(slide, Inches(0.75), Inches(1.4), Inches(4.0), Inches(1.75), "From job_description.docx", [
        "Senior AI Engineer — Founding Team @ Redrob AI (Series A).",
        "Pune/Noida preferred · 5–9 years · production IR over research.",
        "Must-have: embeddings retrieval, vector DB, Python, NDCG/MRR/MAP eval.",
        "Rejects: research-only, LangChain-only, consulting-only, CV without IR.",
    ], body_size=10)

    _card(slide, Inches(0.55), Inches(3.38), Inches(4.4), Inches(1.72), tint=True)
    _rich_bullets(slide, Inches(0.75), Inches(3.5), Inches(4.0), Inches(1.5), "candidate_schema.json", [
        "profile · career_history · skills · redrob_signals (23 fields).",
        "Skills carry proficiency, endorsements, duration_months as trust signals.",
    ], body_size=10)

    _card(slide, Inches(5.1), Inches(1.28), Inches(4.35), Inches(3.82))
    _rich_bullets(slide, Inches(5.3), Inches(1.4), Inches(3.95), Inches(3.55),
                  "23 Redrob Signals (weighted)", [
        "Availability: open_to_work_flag, last_active_date, notice_period_days",
        "Responsiveness: recruiter_response_rate, avg_response_time_hours",
        "Engagement: profile_views, saved_by_recruiters, search_appearance",
        "Quality: skill_assessment_scores, github_activity_score",
        "Trust: verified_email, verified_phone, linkedin_connected",
        "Intent: willing_to_relocate, preferred_work_mode, salary range",
    ], body_size=9.5)


def slide_methodology(prs: Presentation):
    slide = _blank(prs)
    _slide_canvas(slide, "Methodology", 4)
    _slide_title(slide, "Ranking Methodology", "Hybrid v6 — retrieve, score, rank, explain")

    _card(slide, Inches(0.55), Inches(1.28), Inches(4.4), Inches(2.05))
    _rich_bullets(slide, Inches(0.75), Inches(1.4), Inches(4.0), Inches(1.85), "Pipeline", [
        "1  Parse job_description.docx → jd_config.py + JD embedding",
        "2  Load candidates.jsonl (100K) → CandidateIndex",
        "3  Score: hybrid v6 + redrob_signals + honeypot traps",
        "4  Sort ranks 1–100; tie-break candidate_id ascending",
        "5  Export submission.csv — UTF-8, monotonic scores",
    ], body_size=10.5)

    _card(slide, Inches(0.55), Inches(3.48), Inches(4.4), Inches(1.62), tint=True)
    _rich_bullets(slide, Inches(0.75), Inches(3.6), Inches(4.0), Inches(1.4), "Models", [
        "MiniLM-L6-v2 bi-encoder (embeddings.fp16.npz, 100K vectors)",
        "TF-IDF JD overlap · heuristic modifiers · CE OFF",
    ], body_size=10.5)

    _card(slide, Inches(5.1), Inches(1.28), Inches(4.35), Inches(3.82))
    _text(slide, Inches(5.3), Inches(1.42), Inches(3.9), Inches(0.3), "v6 Signal Weights", size=15, bold=True, color=FOREST)
    weights = [
        ("Title alignment", 20), ("Core IR skills", 18), ("Career semantic", 14),
        ("Production pedigree", 12), ("Availability", 12), ("Assessments", 8),
        ("JD overlap", 6), ("Experience", 5), ("Engagement", 5), ("Location", 3),
    ]
    y = 1.78
    for label, pct in weights:
        _text(slide, Inches(5.35), Inches(y), Inches(2.5), Inches(0.2), label, size=10, color=INK)
        _text(slide, Inches(8.05), Inches(y), Inches(0.55), Inches(0.2), f"{pct}%", size=10, bold=True, color=SAGE, align=PP_ALIGN.RIGHT)
        _round_rect(slide, Inches(5.35), Inches(y + 0.2), Inches(3.25), Inches(0.1), CARD_TINT, line=MINT_PALE)
        bar = pct / 20.0 * 3.25
        if bar > 0:
            _round_rect(slide, Inches(5.35), Inches(y + 0.2), Inches(bar), Inches(0.1), LEAF)
        y += 0.34


def slide_explainability(prs: Presentation):
    slide = _blank(prs)
    _slide_canvas(slide, "Explainability", 5)
    _slide_title(slide, "Explainability & Data Validation", "Transparent decisions judges can audit")

    _card(slide, Inches(0.55), Inches(1.28), Inches(4.4), Inches(3.82))
    _rich_bullets(slide, Inches(0.75), Inches(1.4), Inches(4.0), Inches(3.55), "Reasoning Quality", [
        "2 sentences: specific facts + JD link + honest concerns.",
        "Only cites fields present in profile — zero hallucination.",
        "Avoids: empty, identical, or templated reasoning strings.",
        "Rank tone matches position (strong at #5, candid at #95).",
        "Validated: mock_stage4_review — 10/10 rows, 6/6 checks each.",
    ], body_size=10.5)

    _card(slide, Inches(5.1), Inches(1.28), Inches(4.35), Inches(3.82), tint=True)
    _rich_bullets(slide, Inches(5.3), Inches(1.4), Inches(3.95), Inches(3.55), "Trap Defense", [
        "~80 honeypots: impossible tenure, expert skills at 0 months.",
        "Keyword stuffers: weak titles + inflated AI skill lists.",
        "Behavioral twins: same profile, different signals — availability wins.",
        "Template blurbs: recycled descriptions demoted.",
        "sample_submission.csv = format only (ranks honeypots by design).",
        "Our result: 0 honeypots in top-100 (limit: 10%).",
    ], body_size=10)


def slide_workflow(prs: Presentation):
    slide = _blank(prs)
    _slide_canvas(slide, "Workflow", 6)
    _slide_title(slide, "End-to-End Workflow", "From official JD to submission-ready CSV")

    steps = [
        ("01", "JD Parse", "job_description.docx"),
        ("02", "Index", "candidates.jsonl"),
        ("03", "Embed", "MiniLM .npz"),
        ("04", "Score", "Hybrid v6"),
        ("05", "Rank", "Top 100"),
        ("06", "Export", "submission.csv"),
    ]
    x = 0.48
    for num, title, body in steps:
        _round_rect(slide, Inches(x), Inches(1.45), Inches(1.42), Inches(2.35), CARD, line=MINT_PALE)
        _shape(slide, MSO_SHAPE.OVAL, Inches(x + 0.48), Inches(1.58), Inches(0.46), Inches(0.46), SAGE)
        _text(slide, Inches(x + 0.48), Inches(1.63), Inches(0.46), Inches(0.35), num, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(slide, Inches(x + 0.08), Inches(2.2), Inches(1.26), Inches(0.35), title, size=12, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
        _text(slide, Inches(x + 0.06), Inches(2.6), Inches(1.3), Inches(1.0), body, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        if x < 7.8:
            _text(slide, Inches(x + 1.44), Inches(2.15), Inches(0.2), Inches(0.3), "›", size=20, bold=True, color=LEAF, align=PP_ALIGN.CENTER)
        x += 1.56

    _round_rect(slide, Inches(0.55), Inches(4.15), Inches(8.9), Inches(0.88), CARD_TINT, line=MINT)
    _text(slide, Inches(0.8), Inches(4.32), Inches(8.4), Inches(0.55),
          "Reproduce: ./scripts/reproduce_ranking.sh  ·  Byte-verified via verify_submission_artifact.py",
          size=10.5, color=INK)


def slide_architecture(prs: Presentation):
    slide = _blank(prs)
    _slide_canvas(slide, "Architecture", 7)
    _slide_title(slide, "System Architecture", "Graded submission path vs optional demo platform")

    _card(slide, Inches(0.55), Inches(1.28), Inches(8.9), Inches(2.15))
    _text(slide, Inches(0.8), Inches(1.42), Inches(8.4), Inches(0.28), "GRADED PATH", size=11, bold=True, color=LEAF)
    modules = ["rank.py", "redrob_ranker", "features", "embeddings", "honeypot", "availability"]
    x = 0.75
    for i, name in enumerate(modules):
        fill = SAGE if i % 2 == 0 else LEAF
        _round_rect(slide, Inches(x), Inches(1.82), Inches(1.35), Inches(0.52), fill)
        _text(slide, Inches(x), Inches(1.94), Inches(1.35), Inches(0.35), name, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(modules) - 1:
            _text(slide, Inches(x + 1.36), Inches(1.92), Inches(0.2), Inches(0.3), "→", size=14, color=LEAF, align=PP_ALIGN.CENTER)
        x += 1.48
    _text(slide, Inches(0.8), Inches(2.55), Inches(8.4), Inches(0.7),
          "Input: candidates.jsonl + embeddings.fp16.npz  →  Output: submission.csv (100 rows)\n"
          "+ career_blurb.py · assessment.py · rerank.py (cross-encoder OFF)",
          size=9.5, color=MUTED)

    _card(slide, Inches(0.55), Inches(3.62), Inches(8.9), Inches(1.45), tint=True)
    _text(slide, Inches(0.8), Inches(3.76), Inches(8.4), Inches(0.28), "DEMO PATH (not graded)", size=11, bold=True, color=LEAF)
    _text(slide, Inches(0.8), Inches(4.12), Inches(8.4), Inches(0.75),
          "Next.js 14 frontend  →  FastAPI + LangGraph backend  →  dashboard, analytics, AI chat\n"
          "Imports top-100 from submission.csv — does not produce the graded artifact.",
          size=10.5, color=INK)


def slide_results(prs: Presentation):
    slide = _blank(prs)
    _slide_canvas(slide, "Results", 8)
    _slide_title(slide, "Results & Performance", "Evidence-backed outcomes and spec compliance")

    _card(slide, Inches(0.55), Inches(1.28), Inches(4.4), Inches(2.55))
    _text(slide, Inches(0.75), Inches(1.4), Inches(4.0), Inches(0.28), "Top-5 Finalists", size=14, bold=True, color=FOREST)
    tops = [
        ("#1", "CAND_0002025", "0.9900", "Sr AI Engineer @ Apple"),
        ("#2", "CAND_0046064", "0.9324", "Sr NLP Engineer @ Salesforce"),
        ("#3", "CAND_0010685", "0.8659", "NLP Engineer @ Rephrase.ai"),
        ("#4", "CAND_0055905", "0.8510", "Sr ML Engineer @ Flipkart"),
        ("#5", "CAND_0018499", "0.8041", "Sr ML Engineer @ Zomato"),
    ]
    y = 1.75
    for i, (rank, cid, score, role) in enumerate(tops):
        bg = CARD_TINT if i % 2 == 0 else CARD
        _round_rect(slide, Inches(0.75), Inches(y), Inches(4.0), Inches(0.34), bg, line=MINT_PALE)
        _text(slide, Inches(0.88), Inches(y + 0.05), Inches(0.4), Inches(0.22), rank, size=10, bold=True, color=SAGE)
        _text(slide, Inches(1.25), Inches(y + 0.05), Inches(3.35), Inches(0.22), f"{cid}  ·  {score}  ·  {role}", size=9, color=INK)
        y += 0.38

    _card(slide, Inches(5.1), Inches(1.28), Inches(4.35), Inches(2.55), tint=True)
    _rich_bullets(slide, Inches(5.3), Inches(1.4), Inches(3.95), Inches(2.3), "Validation", [
        "100 rows · unique ranks · monotonic scores · 0 honeypots",
        "All IDs verified in official candidates.jsonl",
        "Byte-reproducible artifact: PASS",
        "Reasoning: 100% unique · Stage 4: 10/10 pass",
    ], body_size=10.5)

    _stat_tile(slide, Inches(0.55), Inches(4.05), "~49–60s", "Runtime on 100K", "≪ 5 min limit")
    _stat_tile(slide, Inches(3.25), Inches(4.05), "CPU", "No GPU · No network", "Spec compliant")
    _stat_tile(slide, Inches(5.95), Inches(4.05), "16 GB", "RAM envelope", "~68MB embeddings")


def slide_tech(prs: Presentation):
    slide = _blank(prs)
    _slide_canvas(slide, "Technology", 9)
    _slide_title(slide, "Technologies Used", "Deliberate stack choices for reproducibility and scale")

    _card(slide, Inches(0.55), Inches(1.28), Inches(4.4), Inches(3.82))
    _rich_bullets(slide, Inches(0.75), Inches(1.4), Inches(4.0), Inches(3.55), "Submission Ranker", [
        "Python 3.11 · NumPy · PyYAML — no torch at rank time",
        "all-MiniLM-L6-v2 — committed embeddings.fp16.npz",
        "Docker: docker-compose.ranker.yml for Stage 3 reproduction",
        "pytest: 31 tests passing in challenge/test_ranker.py",
        "AI tools: Cursor + Grok (dev only — no candidate data to LLMs)",
    ], body_size=10.5)

    _card(slide, Inches(5.1), Inches(1.28), Inches(4.35), Inches(3.82), tint=True)
    _rich_bullets(slide, Inches(5.3), Inches(1.4), Inches(3.95), Inches(3.55), "Demo Platform", [
        "Next.js 14 + Tailwind — recruiter command center",
        "FastAPI + LangGraph — 7-agent explainability pipeline",
        "PostgreSQL / SQLite — workspace + top-100 import",
        "HuggingFace Spaces — sample ranking sandbox",
        "Vercel + Render — recruitgpt-x.vercel.app",
    ], body_size=10.5)


def slide_assets(prs: Presentation):
    slide = _blank(prs)
    _slide_canvas(slide, "Assets", 10)
    _slide_title(slide, "Submission Assets", "Complete reproducibility package for judges")

    assets = [
        ("GitHub", "github.com/rahulx2001/recruitgpt-x", "rank.py + challenge/ + bundle validators"),
        ("Ranked CSV", "submission.csv", "submission_spec v4 · UTF-8 · 100 rows"),
        ("HF Sandbox", "huggingface.co/.../recruitgpt-ranker", "§10.5 sample_candidates.json"),
        ("Official Data", "candidates.jsonl + job_description.docx", "sync_challenge_data.sh"),
        ("Reproduce", "./scripts/reproduce_ranking.sh", "CPU ≤5min · byte-verified"),
        ("Metadata", "submission_metadata.yaml", "Team Schadn · Rahul Kumar Singh"),
    ]
    y = 1.3
    for i, (title, url, note) in enumerate(assets):
        _round_rect(slide, Inches(0.55), Inches(y), Inches(8.9), Inches(0.58),
                    CARD_TINT if i % 2 else CARD, line=MINT_PALE)
        _rect(slide, Inches(0.55), Inches(y), Inches(0.08), Inches(0.58), LEAF if i % 2 else SAGE)
        _text(slide, Inches(0.78), Inches(y + 0.07), Inches(1.55), Inches(0.22), title, size=11, bold=True, color=FOREST)
        _text(slide, Inches(2.35), Inches(y + 0.07), Inches(4.8), Inches(0.22), url, size=9.5, color=SAGE)
        _text(slide, Inches(0.78), Inches(y + 0.3), Inches(8.4), Inches(0.2), note, size=8.5, color=MUTED)
        y += 0.63


def slide_closing(prs: Presentation):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CREAM)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), LEAF)
    _shape(slide, MSO_SHAPE.OVAL, Inches(7.5), Inches(-0.8), Inches(3.5), Inches(3.5), MINT_PALE)
    _shape(slide, MSO_SHAPE.OVAL, Inches(-0.6), Inches(4.2), Inches(2.4), Inches(2.4), CARD_TINT)

    _text(slide, Inches(0.9), Inches(1.5), Inches(8.2), Inches(0.75), "Thank You", size=42, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
    _text(slide, Inches(0.9), Inches(2.45), Inches(8.2), Inches(0.55),
          "Schadn  ·  Rahul Kumar Singh", size=18, color=SAGE, align=PP_ALIGN.CENTER)
    _text(slide, Inches(0.9), Inches(3.1), Inches(8.2), Inches(0.45),
          "Ranking candidates like a great recruiter — not a keyword filter.",
          size=13, color=MUTED, align=PP_ALIGN.CENTER)
    _round_rect(slide, Inches(3.1), Inches(3.85), Inches(3.8), Inches(0.55), CARD, line=MINT)
    _text(slide, Inches(3.1), Inches(4.0), Inches(3.8), Inches(0.35),
          "github.com/rahulx2001/recruitgpt-x", size=11, bold=True, color=FOREST, align=PP_ALIGN.CENTER)


def main():
    global _slide_no
    _slide_no = 0
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_solution(prs)
    slide_jd(prs)
    slide_methodology(prs)
    slide_explainability(prs)
    slide_workflow(prs)
    slide_architecture(prs)
    slide_results(prs)
    slide_tech(prs)
    slide_assets(prs)
    slide_closing(prs)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()